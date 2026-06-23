import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from src.models import SlideContent, StyleConfig, BoundingBox
from src.utils import log

# ─── Colour / Luminance Helpers ──────────────────────────────────────────────

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert #RRGGBB to pptx RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _get_luminance(hex_color: str) -> float:
    """Relative luminance (0..1) – copied from style.py to avoid import cycle."""
    hex_color = hex_color.lstrip('#')
    try:
        rgb = tuple(int(hex_color[i:i+2], 16) / 255.0 for i in (0, 2, 4))
    except ValueError:
        return 0.5  # fallback
    a = [c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4 for c in rgb]
    return a[0] * 0.2126 + a[1] * 0.7152 + a[2] * 0.0722

def _is_dark_bg(hex_color: str) -> bool:
    """True if the background luminance is less than 0.5."""
    return _get_luminance(hex_color) < 0.5

# ─── UI Card / Text Block Helpers ────────────────────────────────────────────

def _draw_ui_card(slide, box: BoundingBox, bg_color: RGBColor, is_dark: bool):
    """
    Draws a subtle semi‑transparent card behind content to anchor it.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(box.left - 0.1), Inches(box.top - 0.1),
        Inches(box.width + 0.2), Inches(box.height + 0.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    # Adjust brightness for subtle contrast – 0.8 for dark, -0.2 for light
    shape.fill.fore_color.brightness = 0.8 if is_dark else -0.2
    shape.line.fill.background()   # no border
    return shape

def _render_text_block(slide, box_spec: BoundingBox, text_lines: list[str],
                       style: StyleConfig, color: RGBColor,
                       is_dark: bool, font_size=18, bold=False):
    """
    Renders a block of text with a subtle background card and the given font color.
    """
    if not box_spec or not text_lines:
        return

    # Draw the card first (so text sits on top)
    _draw_ui_card(slide, box_spec, hex_to_rgb(style.bg_color), is_dark)

    # Add the actual text box
    box = slide.shapes.add_textbox(
        Inches(box_spec.left), Inches(box_spec.top),
        Inches(box_spec.width), Inches(box_spec.height)
    )
    tf = box.text_frame
    tf.word_wrap = True

    for i, text in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {text}" if not bold else text
        p.font.name = style.body_font if not bold else style.title_font
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        if i > 0:
            p.space_before = Pt(14)

def _add_visual_placeholder(slide, v_left, v_top, v_width, v_height,
                            is_dark: bool, visual_type: str) -> None:
    """Add a styled placeholder box for missing visual content."""
    ph_box = slide.shapes.add_shape(1, v_left, v_top, v_width, v_height)
    ph_box.fill.solid()
    ph_bg = "#222222" if is_dark else "#F0F0F0"
    ph_text_color = "#888888" if is_dark else "#777777"
    ph_box.fill.fore_color.rgb = hex_to_rgb(ph_bg)
    ph_box.line.color.rgb = hex_to_rgb(ph_text_color)

    tf_ph = ph_box.text_frame
    tf_ph.word_wrap = True
    p_ph = tf_ph.paragraphs[0]
    p_ph.text = f"[ Placeholder for {visual_type.upper()} ]"
    p_ph.alignment = PP_ALIGN.CENTER
    p_ph.font.color.rgb = hex_to_rgb(ph_text_color)

# ─── Main Slide Builder ──────────────────────────────────────────────────────

def _add_slide(prs: Presentation, content: SlideContent, style: StyleConfig) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    W = prs.slide_width

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(style.bg_color)

    is_dark = _is_dark_bg(style.bg_color)

    # ── Use the theme's text_color ──────────────────────────────────
    body_color = hex_to_rgb(style.text_color)
    log("renderer", f"bg={style.bg_color} text_color={style.text_color}")

    # Title accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_to_rgb(style.accent_color)
    accent.line.fill.background()

    # Title
    if content.layout and content.layout.title_box:
        tb = content.layout.title_box
        title_box = slide.shapes.add_textbox(Inches(tb.left), Inches(tb.top),
                                             Inches(tb.width), Inches(tb.height))
    else:
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3),
                                             Inches(9.0), Inches(1.4))

    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.title
    p.font.name = style.title_font
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(style.accent_color)

    # Dynamic content
    if content.layout:
        if content.bullets and content.layout.body_box:
            _render_text_block(slide, content.layout.body_box,
                               content.bullets[:style.max_bullets],
                               style, body_color, is_dark)

        if content.left_column and content.layout.left_box:
            _render_text_block(slide, content.layout.left_box,
                               content.left_column[:style.max_bullets],
                               style, body_color, is_dark, font_size=16)
        if content.right_column and content.layout.right_box:
            _render_text_block(slide, content.layout.right_box,
                               content.right_column[:style.max_bullets],
                               style, body_color, is_dark, font_size=16)

        if content.big_message and content.layout.big_message_box:
            _render_text_block(slide, content.layout.big_message_box,
                               [content.big_message], style,
                               hex_to_rgb(style.accent_color),
                               is_dark, font_size=32, bold=True)

        if content.layout.has_visual and content.layout.visual_box:
            vb = content.layout.visual_box
            v_left, v_top, v_width, v_height = (
                Inches(vb.left), Inches(vb.top),
                Inches(vb.width), Inches(vb.height)
            )
            _draw_ui_card(slide, vb, hex_to_rgb(style.bg_color), is_dark)

            chart_file = f"temp_assets/chart_slide_{content.slide_id}.png"
            if content.visual_hint == "image" and getattr(content, "image_path", None):
                slide.shapes.add_picture(content.image_path, v_left, v_top,
                                         width=v_width, height=v_height)
            elif content.visual_hint == "chart" and os.path.exists(chart_file):
                slide.shapes.add_picture(chart_file, v_left, v_top,
                                         width=v_width, height=v_height)
            else:
                _add_visual_placeholder(slide, v_left, v_top, v_width, v_height,
                                        is_dark, content.visual_hint)

    # Takeaway bar
    if content.takeaway:
        bar = slide.shapes.add_shape(1, Inches(0), Inches(6.6), W, Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(style.accent_color)
        bar.line.fill.background()

        ta_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.65),
                                          Inches(9.3), Inches(0.8))
        tf = ta_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Key takeaway: {content.takeaway}"
        p.font.name = style.body_font
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb("#ffffff")

    if content.speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = content.speaker_notes

# ─── Public Renderer ─────────────────────────────────────────────────────────

def render_pptx(
    slides: list[SlideContent],
    style: StyleConfig,
    output_path: str = "./outputs/presentation.pptx",
) -> str:
    """
    Render the slide contents to a .pptx file using the given style.
    Returns the output file path.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for content in slides:
        _add_slide(prs, content, style)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    # Integrity checks
    check = Presentation(output_path)
    assert len(check.slides) == len(slides), \
        f"Expected {len(slides)} slides, got {len(check.slides)}"

    size = os.path.getsize(output_path)
    assert size > 0, "Output file is empty"

    print(f"[renderer] Saved {len(slides)} slides → {output_path} ({size:,} bytes)")
    return output_path

# ─── Quick Self‑Test ─────────────────────────────────────────────────────────

if __name__ == "__main__":
    sys.path.insert(0, "src")
    from query_compiler import compile_query
    from retrieval import retrieve
    from strategy import generate_strategy
    from planner import generate_slide_plan
    from content import generate_content
    from visual import process_visuals
    from layout import plan_layout
    from style import resolve_style

    raw_query = (
        sys.argv[1] if len(sys.argv) > 1
        else "make a 10 slide technical presentation on the transformer architecture"
    )
    doc_id = sys.argv[2] if len(sys.argv) > 2 else "Attention_is_all_you_Need"

    request    = compile_query(raw_query)
    evidence   = retrieve(request.topic, top_k=15, doc_id=doc_id)
    strategy   = generate_strategy(request, evidence)
    plan       = generate_slide_plan(request, evidence, strategy=strategy)
    slides     = generate_content(plan, evidence, audience=request.audience, strategy=strategy)
    slides     = plan_layout(slides)
    slides     = process_visuals(slides)
    style      = resolve_style(request.style_desc)
    output     = render_pptx(slides, style)

    print(f"\n✅ Done — open {output} in Keynote or PowerPoint")